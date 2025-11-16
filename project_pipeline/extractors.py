# nlp_pipeline/extractors.py
import pdfplumber
from pptx import Presentation
from typing import List, Dict

def extract_text_from_pdf(path: str) -> List[Dict]:
    results = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            results.append({"page": i+1, "text": text})
    return results

def extract_text_from_pptx(path: str) -> List[Dict]:
    prs = Presentation(path)
    results = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
        results.append({"slide": i+1, "text": "\n".join(texts)})
    return results

def get_text_from_file(path: str) -> str:
    if path.lower().endswith(".pdf"):
        pages = extract_text_from_pdf(path)
        return "\n\n".join([p["text"] for p in pages])
    elif path.lower().endswith(".pptx"):
        slides = extract_text_from_pptx(path)
        return "\n\n".join([s["text"] for s in slides])
    else:
        raise ValueError("Unsupported file format. Use PDF or PPTX.")
